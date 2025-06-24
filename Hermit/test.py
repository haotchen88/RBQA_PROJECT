import streamlit as st
import time
from datetime import datetime
import os
import pandas as pd
import json
from docx import Document
from PyPDF2 import PdfReader
import pptx
import tempfile
import hashlib
import numpy as np
from sklearn.metrics.pairwise import cosine_similarity
import re
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
from sentence_transformers import SentenceTransformer

# 下载NLTK资源（首次运行时需要）
try:
    nltk.data.find('tokenizers/punkt')
    nltk.data.find('corpora/stopwords')
except LookupError:
    nltk.download('punkt')
    nltk.download('stopwords')


# 初始化会话状态
def init_session():
    if "conversation" not in st.session_state:
        st.session_state.conversation = []
    if "uploaded_files" not in st.session_state:
        st.session_state.uploaded_files = []
    if "knowledge_base" not in st.session_state:
        st.session_state.knowledge_base = []
    if "deleted_files" not in st.session_state:
        st.session_state.deleted_files = []
    if "embedding_model" not in st.session_state:
        st.session_state.embedding_model = None
    if "knowledge_vectors" not in st.session_state:
        st.session_state.knowledge_vectors = []


# 加载嵌入模型（缓存避免重复加载）
@st.cache_resource
def load_embedding_model():
    return SentenceTransformer('paraphrase-multilingual-MiniLM-L12-v2')


# 生成文件唯一ID
def generate_file_id(file_content):
    return hashlib.md5(file_content).hexdigest()[:8]


# 文件解析函数
def parse_file(file):
    content = ""
    file_type = file.name.split(".")[-1].lower()

    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_type}") as tmp:
            tmp.write(file.getvalue())
            tmp_path = tmp.name

        if file_type == "txt":
            with open(tmp_path, "r", encoding="utf-8") as f:
                content = f.read()

        elif file_type == "pdf":
            reader = PdfReader(tmp_path)
            for page in reader.pages:
                content += page.extract_text() + "\n"

        elif file_type == "docx":
            doc = Document(tmp_path)
            for para in doc.paragraphs:
                content += para.text + "\n"

        elif file_type == "pptx":
            prs = pptx.Presentation(tmp_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"

        os.unlink(tmp_path)
        return content

    except Exception as e:
        st.error(f"解析文件时出错: {str(e)}")
        return None


# 文本预处理函数
def preprocess_text(text):
    # 转换为小写
    text = text.lower()
    # 移除特殊字符和数字
    text = re.sub(r'[^a-zA-Z\u4e00-\u9fff\s]', '', text)
    # 分词
    tokens = word_tokenize(text)
    # 移除停用词
    stop_words = set(stopwords.words('english') + stopwords.words('chinese'))
    tokens = [word for word in tokens if word not in stop_words]
    return " ".join(tokens)


# 将文档内容分割为适合处理的片段
def split_into_chunks(text, chunk_size=200):
    chunks = []
    words = text.split()

    for i in range(0, len(words), chunk_size):
        chunks.append(" ".join(words[i:i + chunk_size]))

    return chunks


# 语义分析函数
def semantic_analysis(question):
    # 1. 意图识别
    intent = "信息查询"  # 默认意图

    if any(word in question.lower() for word in ["如何", "怎样", "步骤"]):
        intent = "操作指导"
    elif any(word in question.lower() for word in ["为什么", "原因", "为何"]):
        intent = "原因解释"
    elif any(word in question.lower() for word in ["比较", "对比", "vs"]):
        intent = "比较分析"
    elif any(word in question.lower() for word in ["推荐", "建议", "应该"]):
        intent = "推荐建议"

    # 2. 关键实体提取（简化版）
    # 使用正则表达式提取可能的实体
    entities = re.findall(r'[\u4e00-\u9fff]{2,}|[a-zA-Z]{3,}', question)

    # 3. 生成语义向量
    model = st.session_state.embedding_model
    embedding = model.encode([question])[0]

    return {
        "intent": intent,
        "entities": entities,
        "embedding": embedding
    }


# 删除文件处理
def delete_file(file_id):
    # 从已上传文件中删除
    st.session_state.uploaded_files = [f for f in st.session_state.uploaded_files if f['id'] != file_id]

    # 添加到删除的文件列表（用于恢复）
    deleted_file = next((f for f in st.session_state.deleted_files if f['id'] == file_id), None)
    if not deleted_file:
        file_to_delete = next((f for f in st.session_state.uploaded_files if f['id'] == file_id), None)
        if file_to_delete:
            st.session_state.deleted_files.append({
                **file_to_delete,
                'deleted_time': datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            })

    # 从知识库中删除相关片段
    st.session_state.knowledge_base = [kb for kb in st.session_state.knowledge_base if kb['source_id'] != file_id]

    # 更新向量存储
    update_vector_store()


# 恢复已删除文件
def restore_file(file_id):
    # 从删除列表中恢复
    file_to_restore = next((f for f in st.session_state.deleted_files if f['id'] == file_id), None)
    if file_to_restore:
        st.session_state.uploaded_files.append({
            'id': file_to_restore['id'],
            'name': file_to_restore['name'],
            'type': file_to_restore['type'],
            'content': file_to_restore['content'],
            'upload_time': file_to_restore['upload_time'],
            'tags': file_to_restore['tags'] if 'tags' in file_to_restore else []
        })

        # 从删除列表中移除
        st.session_state.deleted_files = [f for f in st.session_state.deleted_files if f['id'] != file_id]

        # 重新添加到知识库
        chunks = split_into_chunks(file_to_restore['content'])
        for i, chunk in enumerate(chunks):
            st.session_state.knowledge_base.append({
                "source": file_to_restore['name'],
                "source_id": file_to_restore['id'],
                "content": chunk,
                "type": file_to_restore['type'].split("/")[-1]
            })

        # 更新向量存储
        update_vector_store()
        return True
    return False


# 标记文件功能
def toggle_file_tag(file_id, tag):
    for file in st.session_state.uploaded_files:
        if file['id'] == file_id:
            if 'tags' not in file:
                file['tags'] = []

            if tag in file['tags']:
                file['tags'].remove(tag)
            else:
                file['tags'].append(tag)
            return True
    return False


# 更新向量存储
def update_vector_store():
    if st.session_state.embedding_model and st.session_state.knowledge_base:
        model = st.session_state.embedding_model
        contents = [kb["content"] for kb in st.session_state.knowledge_base]
        st.session_state.knowledge_vectors = model.encode(contents)


# 知识库管理界面（添加文件删除和标记功能）
def knowledge_base_section():
    st.header("📚 知识库构建与管理")

    # 文件上传区域
    uploaded_files = st.file_uploader(
        "上传知识文档 (支持多种格式)",
        type=["txt", "pdf", "docx", "pptx", "md"],
        accept_multiple_files=True
    )

    if uploaded_files:
        # 确保嵌入模型已加载
        if not st.session_state.embedding_model:
            with st.spinner("加载语义模型..."):
                st.session_state.embedding_model = load_embedding_model()

        # 处理新上传的文件
        for file in uploaded_files:
            # 检查是否已上传过相同内容的文件
            file_id = generate_file_id(file.getvalue())
            existing_file = next((f for f in st.session_state.uploaded_files if f['id'] == file_id), None)

            if not existing_file:
                with st.spinner(f"解析文件: {file.name}..."):
                    content = parse_file(file)

                    if content:
                        st.session_state.uploaded_files.append({
                            "id": file_id,
                            "name": file.name,
                            "type": file.type,
                            "content": content,
                            "upload_time": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                            "tags": ["新上传"]  # 默认标记
                        })

                        # 分割内容为知识片段
                        chunks = split_into_chunks(content)
                        for i, chunk in enumerate(chunks):
                            st.session_state.knowledge_base.append({
                                "source": file.name,
                                "source_id": file_id,
                                "content": chunk,
                                "type": file.type.split("/")[-1]
                            })

        # 更新向量存储
        update_vector_store()
        st.rerun()

    # 显示上传文件列表
    st.subheader("文档管理")
    st.info("使用以下表格管理您的文档：")

    files_df = pd.DataFrame([
        {
            "ID": f["id"],
            "文件名": f["name"],
            "类型": f["type"],
            "大小": f"{len(f['content']) // 1024} KB",
            "上传时间": f["upload_time"],
            "标记": ", ".join(f.get("tags", [])) if "tags" in f else "",
        }
        for f in st.session_state.uploaded_files
    ])

    # 如果没有任何文件显示提示信息
    if files_df.empty:
        st.info("暂无文档，请上传文档")
    else:
        # 显示数据表格
        st.dataframe(files_df.set_index('ID'), use_container_width=True)

        # 文件管理功能区
        with st.expander("📌 文件操作", expanded=True):
            col1, col2, col3, col4 = st.columns([0.4, 0.2, 0.2, 0.2])

            # 文件选择
            file_ids = list(files_df["ID"])
            selected_file_id = col1.selectbox("选择文件", file_ids, format_func=lambda id:
            files_df.loc[files_df["ID"] == id, "文件名"].values[0])

            # 标记管理
            tags = ["重要", "待审核", "存档", "参考"]
            selected_tag = col2.selectbox("添加/移除标记", tags)

            # 按钮操作
            if col3.button("应用标记", key="tag_btn", use_container_width=True):
                toggle_file_tag(selected_file_id, selected_tag)
                st.rerun()

            if col4.button("删除文件", key="delete_btn", type="primary", use_container_width=True):
                delete_file(selected_file_id)
                st.rerun()

        # 回收站管理
        if st.session_state.deleted_files:
            with st.expander("🗑️ 回收站管理", expanded=True):
                deleted_df = pd.DataFrame([
                    {
                        "ID": f["id"],
                        "文件名": f["name"],
                        "类型": f["type"],
                        "删除时间": f["deleted_time"],
                    }
                    for f in st.session_state.deleted_files
                ])

                st.dataframe(deleted_df.set_index('ID'), use_container_width=True)

                restore_id = st.selectbox("选择要恢复的文件", deleted_df["ID"], format_func=lambda id:
                deleted_df.loc[deleted_df["ID"] == id, "文件名"].values[0])

                if st.button("恢复文件", use_container_width=True):
                    if restore_file(restore_id):
                        st.success(f"文件已恢复")
                        st.rerun()
                    else:
                        st.error("恢复文件失败")

                if st.button("清空回收站", type="primary", use_container_width=True):
                    st.session_state.deleted_files = []
                    st.success("回收站已清空")
                    st.rerun()

        # 查看知识库内容
        with st.expander("🔍 查看知识库内容", expanded=False):
            if not st.session_state.knowledge_base:
                st.info("知识库为空")
            else:
                # 分组显示按文件
                source_files = set(kb['source_id'] for kb in st.session_state.knowledge_base)
                for src_id in list(source_files)[:3]:  # 最多显示前3个文件的内容
                    source_name = next(
                        kb['source'] for kb in st.session_state.knowledge_base if kb['source_id'] == src_id)
                    st.subheader(f"来源: {source_name}")

                    # 显示该文件的前3个片段
                    for kb in [k for k in st.session_state.knowledge_base if k['source_id'] == src_id][:3]:
                        with st.expander(f"知识片段 {kb['content'][:30]}...", expanded=False):
                            st.markdown(kb["content"])
                    st.divider()
                if len(source_files) > 3:
                    st.info(f"已显示3个文件的内容，共{len(source_files)}个文件")


# 问答界面（添加语义理解功能）
def qa_interface():
    st.header("💬 智能问答系统 (语义理解版)")

    # 显示历史对话
    if st.session_state.conversation:
        st.subheader("对话历史")
        for i, message in enumerate(st.session_state.conversation):
            role = "🕵️‍♂️ 用户" if i % 2 == 0 else "🤖 系统"
            with st.chat_message(name=role):
                st.write(message)
            if i < len(st.session_state.conversation) - 1:
                st.divider()

    # 用户输入问题
    question = st.chat_input("请输入您的问题...")

    if question:
        # 确保嵌入模型已加载
        if not st.session_state.embedding_model:
            with st.spinner("加载语义模型..."):
                st.session_state.embedding_model = load_embedding_model()

        # 添加到对话历史
        st.session_state.conversation.append(f"用户: {question}")

        with st.chat_message(name="🕵️‍♂️ 用户"):
            st.write(question)

        # 语义分析处理
        with st.spinner("正在分析问题语义..."):
            semantic_info = semantic_analysis(question)
            intent = semantic_info["intent"]
            entities = semantic_info["entities"]
            question_embedding = semantic_info["embedding"]

        # 语义检索过程
        with st.spinner("正在检索相关知识..."):
            # 为知识库片段生成嵌入向量（如果尚未生成）
            if st.session_state.knowledge_vectors.size == 0 and st.session_state.knowledge_base:
                model = st.session_state.embedding_model
                contents = [kb["content"] for kb in st.session_state.knowledge_base]
                st.session_state.knowledge_vectors = model.encode(contents)

            # 计算相似度
            similarities = []
            for idx, kb_vector in enumerate(st.session_state.knowledge_vectors):
                similarity = cosine_similarity([question_embedding], [kb_vector])[0][0]
                similarities.append((idx, similarity))

            # 按相似度排序
            similarities.sort(key=lambda x: x[1], reverse=True)

            # 获取最相关的前3个片段
            top_matches = []
            for idx, score in similarities[:3]:
                if score > 0.3:  # 设置相似度阈值
                    top_matches.append(st.session_state.knowledge_base[idx])

        # 生成答案
        with st.spinner("正在生成答案..."):
            if top_matches:
                answer = f"### 问题分析\n"
                answer += f"- **意图识别**: {intent}\n"
                if entities:
                    answer += f"- **关键实体**: {', '.join(entities)}\n"

                answer += f"\n### 根据相关知识为您解答\n"
                for i, match in enumerate(top_matches):
                    answer += f"**来源 {i + 1}** ({match['source']}):\n"
                    answer += f"> {match['content'][:200]}...\n\n"

                # 根据不同意图添加额外信息
                if intent == "比较分析":
                    answer += "\n> *提示：如需详细比较分析，请提供更多具体信息*"
                elif intent == "推荐建议":
                    answer += "\n> *提示：以上为基于知识库的推荐，仅供参考*"
                elif intent == "操作指导":
                    answer += "\n> *提示：如需更详细的操作步骤指南，请补充具体场景信息*"
            else:
                answer = "在知识库中未找到相关信息。请尝试重新表述您的问题或上传更多相关文档。"

            # 添加到对话历史
            st.session_state.conversation.append(f"系统: {answer}")

            # 显示答案
            with st.chat_message(name="🤖 系统"):
                st.markdown(answer)

            # 调试信息
            with st.expander("🔍 语义分析详情", expanded=False):
                st.json({
                    "问题意图": intent,
                    "识别实体": entities,
                    "匹配片段数": len(top_matches),
                    "最高相似度": similarities[0][1] if similarities else 0
                })


# 主界面
def main():
    st.set_page_config(
        page_title="智能问答系统",
        page_icon="🤖",
        layout="wide",
        initial_sidebar_state="expanded"
    )

    st.title("📚 基于语义理解的智能问答系统")
    st.caption("知识库构建、管理及智能问答平台 | 支持文档处理与语义分析")

    init_session()

    # 创建侧边栏导航
    with st.sidebar:
        st.header("🔍 导航菜单")
        page = st.radio("选择功能", ["知识库管理", "智能问答"], horizontal=True)
        st.divider()

        # 系统状态概览
        st.subheader("📊 系统概览")
        col1, col2, col3 = st.columns(3)
        col1.metric("知识文档", len(st.session_state.uploaded_files))
        col2.metric("知识片段", len(st.session_state.knowledge_base))
        col3.metric("回收站", len(st.session_state.deleted_files))
        st.divider()

        st.info("""
        **系统功能：**
        1. 知识库构建与管理
        2. 文档解析与处理
        3. 语义分析与理解
        4. 智能问答系统
        5. 文件标记与回收
        """)

        # 数据管理选项
        if st.button("清空所有数据", use_container_width=True, type="secondary"):
            st.session_state.uploaded_files = []
            st.session_state.knowledge_base = []
            st.session_state.conversation = []
            st.session_state.knowledge_vectors = []
            st.rerun()

    # 根据选择显示对应页面
    if page == "知识库管理":
        knowledge_base_section()
    else:
        qa_interface()

    # 调试信息
    with st.expander("🛠️ 调试信息", expanded=False):
        st.json({
            "文件上传数": len(st.session_state.uploaded_files),
            "知识片段数": len(st.session_state.knowledge_base),
            "向量存储数": len(st.session_state.knowledge_vectors),
            "删除文件数": len(st.session_state.deleted_files),
            "对话轮次": len(st.session_state.conversation) // 2
        })


if __name__ == "__main__":
    main()