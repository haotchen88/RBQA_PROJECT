import streamlit as st
import time
from datetime import datetime
import os
import pandas as pd
from docx import Document
from PyPDF2 import PdfReader
import pptx
import tempfile
import hashlib
from langchain.vectorstores import Chroma
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
import openai


# ==================== 新增开始 ====================
# DeepSeek AI 接口配置
# 生成文件唯一ID（添加在 init_session() 函数之后）
def generate_file_id(file_content):
    """通过文件内容生成唯一哈希ID"""
    return hashlib.md5(file_content).hexdigest()[:8]

# 文件解析函数（添加在 generate_file_id() 之后）
def parse_file(file):
    """解析上传的文件内容"""
    content = ""
    file_type = file.name.split(".")[-1].lower()

    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_type}") as tmp:
            tmp.write(file.getvalue())
            tmp_path = tmp.name

        if file_type == "txt":
            with open(tmp_path, "r", encoding="utf-8") as f:
                content = f.read()

        elif file_type == "pdf":
            reader = PdfReader(tmp_path)
            for page in reader.pages:
                content += page.extract_text() + "\n"

        elif file_type == "docx":
            doc = Document(tmp_path)
            for para in doc.paragraphs:
                content += para.text + "\n"

        elif file_type == "pptx":
            prs = pptx.Presentation(tmp_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"

        os.unlink(tmp_path)
        return content

    except Exception as e:
        st.error(f"解析文件时出错: {str(e)}")
        return None
def setup_deepseek():
    openai.api_key = "sk-g40Ua40lLiQhMcEN1b710a5d63E14bD89921Ed47D8B371Fb"  # 从secrets获取，不存在则为空
    openai.base_url = "https://api.gpt.ge/v1/"


def ask_ai(question, model="deepseek-chat"):
    """调用DeepSeek AI接口"""
    try:
        completion = openai.chat.completions.create(
            model=model,
            messages=[
                {
                    "role": "system",
                    "content": "你是一个帮助进行文献管理的AI问答系统。根据提供的上下文，用精炼而科学的语言回答问题"
                },
                {"role": "user", "content": question},
            ],
        )
        return completion.choices[0].message.content.strip()
    except Exception as e:
        st.error(f"AI接口调用失败: {str(e)}")
        return "无法获取AI回答，请检查API配置"


# ==================== 新增结束 ====================

# 初始化会话状态
def init_session():
    if "conversation" not in st.session_state:
        st.session_state.conversation = []
    if "uploaded_files" not in st.session_state:
        st.session_state.uploaded_files = []
    if "knowledge_base" not in st.session_state:
        st.session_state.knowledge_base = []
    if "deleted_files" not in st.session_state:
        st.session_state.deleted_files = []
    if "vector_db" not in st.session_state:
        st.session_state.embedding_model = HuggingFaceEmbeddings(
            model_name="GanymedeNil/text2vec-large-chinese",
            model_kwargs={'device': 'cpu'}
        )
        st.session_state.vector_db = Chroma(
            embedding_function=st.session_state.embedding_model,
            persist_directory="./chroma_db"
        )
    if "text_splitter" not in st.session_state:
        st.session_state.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )


# [保留您原有的 parse_file、generate_file_id、delete_file 等辅助函数...]

# 知识库管理界面
def knowledge_base_section():
    st.header("📚 知识库构建与管理")
    uploaded_files = st.file_uploader(
        "上传知识文档 (支持PDF/DOCX/PPTX/TXT)",
        type=["pdf", "docx", "pptx", "txt"],
        accept_multiple_files=True
    )

    if uploaded_files:
        for file in uploaded_files:
            file_id = generate_file_id(file.getvalue())
            if not next((f for f in st.session_state.uploaded_files if f['id'] == file_id), None):
                content = parse_file(file)
                if content:
                    chunks = st.session_state.text_splitter.split_text(content)

                    # 存储到向量数据库
                    st.session_state.vector_db.add_texts(
                        texts=chunks,
                        metadatas=[{
                            "source": file.name,
                            "source_id": file_id,
                            "type": file.type.split("/")[-1],
                            "upload_time": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                        } for _ in chunks]
                    )

                    # 保留到知识库
                    st.session_state.uploaded_files.append({
                        "id": file_id,
                        "name": file.name,
                        "type": file.type,
                        "content": content,
                        "upload_time": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                        "tags": ["新上传"]
                    })

    # [保留您原有的文件列表展示、删除/恢复等UI逻辑...]


# ==================== 修改开始 ====================
# 增强版问答界面（集成DeepSeek）
def qa_interface():
    st.header("💬 智能问答系统")

    # 显示对话历史
    if st.session_state.conversation:
        for msg in st.session_state.conversation:
            role = "user" if msg.startswith("用户:") else "assistant"
            with st.chat_message(role):
                st.write(msg.split(":", 1)[1].strip())

    # 用户提问处理
    if question := st.chat_input("请输入您的问题..."):
        st.session_state.conversation.append(f"用户: {question}")

        # 1. 向量检索
        docs = st.session_state.vector_db.similarity_search(question, k=3)

        # 2. 构建科学问答提示词
        context = "\n".join([
            f"【文献 {i + 1}】{doc.metadata['source']}\n{doc.page_content}\n"
            for i, doc in enumerate(docs)
        ])

        prompt = f"""根据以下文献内容回答问题：
{context}
问题：{question}
要求：
1. 回答需引用文献（例：【文献1】）
2. 保持学术严谨性
3. 如无相关信息请说明
回答："""

        # 3. 调用DeepSeek生成
        with st.chat_message("assistant"):
            with st.spinner("正在生成回答..."):
                answer = ask_ai(prompt)
                st.write(answer)
                st.session_state.conversation.append(f"系统: {answer}")

            # 显示参考文献
            with st.expander("📚 参考文档", expanded=False):
                for i, doc in enumerate(docs, 1):
                    st.caption(f"【文献{i}】{doc.metadata['source']}")
                    st.text(doc.page_content[:200] + "...")


# ==================== 修改结束 ====================

# 主界面
def main():
    st.set_page_config(
        page_title="智能文献问答系统",
        page_icon="📚",
        layout="wide"
    )

    setup_deepseek()  # 初始化DeepSeek配置
    init_session()

    # 侧边栏
    with st.sidebar:
        st.header("导航菜单")
        page = st.radio("功能选择", ["知识库管理", "智能问答"])

        st.divider()
        st.info("""
        **使用说明：**
        1. 上传文献到知识库
        2. 在问答页提问
        3. 系统将结合文献回答
        """)

    if page == "知识库管理":
        knowledge_base_section()
    else:
        qa_interface()


if __name__ == "__main__":
    main()